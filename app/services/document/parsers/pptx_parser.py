import logging
from pptx import Presentation
from .base_parser import BaseParser, ParseResult

logger = logging.getLogger(__name__)


class PptxParser(BaseParser):
    """PPT 文档解析器，使用 python-pptx"""

    def parse(self, file_path: str) -> ParseResult:
        prs = Presentation(file_path)
        slides_text = []

        for slide_num, slide in enumerate(prs.slides, 1):
            slide_parts = [f'--- Slide {slide_num} ---']
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for paragraph in shape.text_frame.paragraphs:
                        text = paragraph.text.strip()
                        if text:
                            slide_parts.append(text)
                if shape.has_table:
                    table = shape.table
                    for row in table.rows:
                        row_text = [cell.text.strip() for cell in row.cells if cell.text.strip()]
                        if row_text:
                            slide_parts.append(' | '.join(row_text))

            if len(slide_parts) > 1:
                slides_text.append('\n'.join(slide_parts))

        content = '\n\n'.join(slides_text)

        # 提取元数据
        core_props = prs.core_properties
        title = core_props.title or ''
        if not title:
            import os
            title = os.path.splitext(os.path.basename(file_path))[0]

        return ParseResult(
            content=content.strip(),
            title=title,
            metadata={
                'author': core_props.author or '',
                'slide_count': len(prs.slides),
                'format': 'pptx',
            }
        )

    def supported_extensions(self) -> list:
        return ['.pptx', '.ppt']
